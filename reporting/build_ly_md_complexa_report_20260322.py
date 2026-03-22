from __future__ import annotations

import importlib.util
import sys
import textwrap
from pathlib import Path

import matplotlib
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import seaborn as sns
from Bio.PDB import PDBParser
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

matplotlib.use("Agg")

ROOT = Path(r"D:\Project\cw\ly")
OLD_PPT = ROOT / "LY_MD_Complexa_Report_20260321.pptx"
OLD_ASSETS = ROOT / "LY_MD_Complexa_Report_assets"
ASSET_DIR = ROOT / "LY_MD_Complexa_Report_20260322_cn_assets"
PPTX_PATH = ROOT / "LY_MD_Complexa_Report_20260322_cn_v2.pptx"
SUMMARY_PATH = ROOT / "LY_MD_Complexa_Report_20260322_cn_v2_summary.txt"

CSV_404 = ROOT / "Complexa_B404_B405_all_pass" / "all_pass_sequences.csv"
CSV_404_FAMILY = ROOT / "Complexa_B404_B405_all_pass" / "family_summary.csv"
CSV_348 = ROOT / "Complexa_B348_all_pass" / "all_pass_sequences.csv"
CSV_348_FAMILY = ROOT / "Complexa_B348_all_pass" / "family_summary.csv"
CSV_COMBO = ROOT / "Joint_peptide_combo_screen_all27" / "all_combo_metrics_annotated.csv"
CSV_SHORTLIST = ROOT / "Joint_peptide_combo_screen_all27" / "experimental_shortlist.csv"


def load_base_module():
    path = Path(__file__).with_name("build_ly_md_complexa_report.py")
    spec = importlib.util.spec_from_file_location("ly_report_base", path)
    module = importlib.util.module_from_spec(spec)
    sys.modules[spec.name] = module
    assert spec.loader is not None
    spec.loader.exec_module(module)
    return module


base = load_base_module()
COLORS = dict(base.COLORS)
FONT_FAMILY = base.FONT_FAMILY

sns.set_style("whitegrid")
plt.rcParams["font.family"] = FONT_FAMILY
plt.rcParams["font.sans-serif"] = [FONT_FAMILY, "SimHei", "DejaVu Sans"]
plt.rcParams["axes.unicode_minus"] = False


def save_fig(fig: plt.Figure, name: str, dpi: int = 220) -> Path:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    path = ASSET_DIR / name
    fig.savefig(path, dpi=dpi, bbox_inches="tight")
    plt.close(fig)
    return path


def rgb_color(value: str) -> RGBColor:
    token = {"white": "FFFFFF", "black": "000000"}.get(value.lower(), value).lstrip("#")
    return RGBColor.from_string(token)


def add_metric_box(slide, left, top, width, height, value: str, label: str, fill: str, line: str | None = None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb_color(fill)
    shape.line.color.rgb = rgb_color(line or fill)
    base.add_textbox(slide, left + Inches(0.12), top + Inches(0.10), width - Inches(0.2), Inches(0.34), value, font_size=24, bold=True, color="white")
    base.add_textbox(slide, left + Inches(0.12), top + Inches(0.48), width - Inches(0.2), Inches(0.28), label, font_size=10, color="white")


def add_source_footer(slide, text: str):
    base.add_textbox(slide, Inches(0.7), Inches(7.02), Inches(12.0), Inches(0.2), text, font_size=8, color=COLORS["gray"], align=PP_ALIGN.RIGHT)


def read_tables():
    df404 = pd.read_csv(CSV_404)
    fam404 = pd.read_csv(CSV_404_FAMILY)
    df348 = pd.read_csv(CSV_348)
    fam348 = pd.read_csv(CSV_348_FAMILY)
    combo = pd.read_csv(CSV_COMBO)
    shortlist = pd.read_csv(CSV_SHORTLIST)
    return df404, fam404, df348, fam348, combo, shortlist


def compute_md_stats():
    rmsd_complex = base.read_xvg(base.MD_DIR / "rmsd_complex.xvg")
    rmsd_a = base.read_xvg(base.MD_DIR / "rmsd_chain_A.xvg")
    rmsd_b = base.read_xvg(base.MD_DIR / "rmsd_chain_B.xvg")
    mindist = base.read_xvg(base.MD_DIR / "mindist.xvg")
    contacts = base.read_xvg(base.MD_DIR / "num_contact.xvg")
    _, hbonds, _ = base.read_xvg_two_series(base.MD_DIR / "hbnum.xvg")
    records = base.parse_hbond_records(base.PLOT_DIR / "hbonds_corrected.pml")
    _, residue_df = base.summarize_hbond_records(records)
    return {
        "rmsd_complex": base.compute_stats(rmsd_complex, base.ANALYSIS_START_NS, base.ANALYSIS_END_NS),
        "rmsd_a": base.compute_stats(rmsd_a, base.ANALYSIS_START_NS, base.ANALYSIS_END_NS),
        "rmsd_b": base.compute_stats(rmsd_b, base.ANALYSIS_START_NS, base.ANALYSIS_END_NS),
        "mindist": base.compute_stats(mindist, base.ANALYSIS_START_NS, base.ANALYSIS_END_NS),
        "contacts": base.compute_stats(contacts, base.ANALYSIS_START_NS, base.ANALYSIS_END_NS),
        "hbonds_mean": float(np.mean(hbonds)),
        "hbonds_std": float(np.std(hbonds)),
        "hotspots": [int(x) for x in residue_df["residue"].head(2).tolist()],
    }


def wrap_label(text: str, width: int = 14) -> str:
    return "\n".join(textwrap.wrap(str(text), width=width, break_long_words=False))


def create_404405_figure(df: pd.DataFrame, families: pd.DataFrame) -> Path:
    fig = plt.figure(figsize=(12.4, 6.8))
    gs = fig.add_gridspec(2, 2, hspace=0.35, wspace=0.28)

    ax1 = fig.add_subplot(gs[0, 0])
    sc = ax1.scatter(df["binder_plddt"], df["i_ptm"], c=df["binder_length"], s=110, cmap="YlGnBu", edgecolor="white", linewidth=1.0)
    top = df.sort_values("secondary_score", ascending=False).head(4)
    for _, row in top.iterrows():
        ax1.text(row["binder_plddt"] + 0.002, row["i_ptm"] + 0.002, row["binder_sequence"], fontsize=8, color=COLORS["ink"])
    ax1.set_title("404/405 候选肽置信度分布", fontsize=13)
    ax1.set_xlabel("肽链 pLDDT")
    ax1.set_ylabel("iPTM")
    cbar = fig.colorbar(sc, ax=ax1, fraction=0.048, pad=0.02)
    cbar.set_label("长度（aa）")

    ax2 = fig.add_subplot(gs[0, 1])
    fam_plot = families.copy().sort_values(["family_size", "representative_score"], ascending=[False, False])
    ax2.bar(fam_plot["family_id"], fam_plot["family_size"], color=["#2A9D8F" if x == "Family A" else "#8C98A4" for x in fam_plot["family_id"]])
    for i, row in enumerate(fam_plot.itertuples(index=False)):
        ax2.text(i, row.family_size + 0.15, str(row.family_size), ha="center", va="bottom", fontsize=9)
    ax2.set_title("序列家族组成", fontsize=13)
    ax2.set_ylabel("家族大小")
    ax2.set_ylim(0, max(fam_plot["family_size"]) + 2)

    ax3 = fig.add_subplot(gs[1, 0])
    length_counts = df["binder_length"].value_counts().sort_index()
    ax3.bar(length_counts.index.astype(str), length_counts.values, color="#4C78A8")
    ax3.set_title("长度分布", fontsize=13)
    ax3.set_xlabel("多肽长度（aa）")
    ax3.set_ylabel("数量")

    ax4 = fig.add_subplot(gs[1, 1])
    top_seq = df.sort_values("secondary_score", ascending=False).head(8).iloc[::-1]
    ax4.barh(top_seq["binder_sequence"], top_seq["secondary_score"], color="#D8A329")
    ax4.set_title("Top 8 候选肽", fontsize=13)
    ax4.set_xlabel("二次筛选得分")

    fig.suptitle("404/405 位点扩样结果：16 条 strict-pass 候选，形成 4 个家族", fontsize=16, fontweight="bold", color=COLORS["navy"])
    return save_fig(fig, "14_404405_design_update.png")


def create_348_figure(df: pd.DataFrame, families: pd.DataFrame) -> Path:
    fig = plt.figure(figsize=(12.4, 6.8))
    gs = fig.add_gridspec(2, 2, hspace=0.35, wspace=0.28)

    ax1 = fig.add_subplot(gs[0, 0])
    sc = ax1.scatter(df["binder_plddt"], df["i_ptm"], c=df["binder_length"], s=110, cmap="YlOrRd", edgecolor="white", linewidth=1.0)
    top = df.sort_values("secondary_score", ascending=False).head(5)
    for _, row in top.iterrows():
        ax1.text(row["binder_plddt"] + 0.002, row["i_ptm"] + 0.002, row["binder_sequence"], fontsize=8, color=COLORS["ink"])
    ax1.set_title("348 位点候选肽置信度分布", fontsize=13)
    ax1.set_xlabel("肽链 pLDDT")
    ax1.set_ylabel("iPTM")
    cbar = fig.colorbar(sc, ax=ax1, fraction=0.048, pad=0.02)
    cbar.set_label("长度（aa）")

    ax2 = fig.add_subplot(gs[0, 1])
    fam_plot = families.copy().sort_values(["family_size", "representative_score"], ascending=[False, False])
    ax2.bar(fam_plot["family_id"], fam_plot["family_size"], color=["#E07A5F" if x == "Family F" else "#8C98A4" for x in fam_plot["family_id"]])
    for i, row in enumerate(fam_plot.itertuples(index=False)):
        ax2.text(i, row.family_size + 0.12, str(row.family_size), ha="center", va="bottom", fontsize=9)
    ax2.set_title("348 位点序列多样性", fontsize=13)
    ax2.set_ylabel("家族大小")
    ax2.set_ylim(0, max(fam_plot["family_size"]) + 2)

    ax3 = fig.add_subplot(gs[1, 0])
    length_counts = df["binder_length"].value_counts().sort_index()
    ax3.bar(length_counts.index.astype(str), length_counts.values, color="#C8553D")
    ax3.set_title("长度分布", fontsize=13)
    ax3.set_xlabel("多肽长度（aa）")
    ax3.set_ylabel("数量")

    ax4 = fig.add_subplot(gs[1, 1])
    fam_show = fam_plot[["family_id", "representative_sequence", "family_size"]].copy()
    ax4.axis("off")
    ax4.set_title("代表序列", fontsize=13, pad=10)
    y = 0.95
    for row in fam_show.itertuples(index=False):
        ax4.text(0.02, y, f"{row.family_id} ({row.family_size})", fontsize=10, fontweight="bold", color=COLORS["navy"], transform=ax4.transAxes)
        ax4.text(0.38, y, row.representative_sequence, fontsize=10, color=COLORS["ink"], transform=ax4.transAxes)
        y -= 0.1

    fig.suptitle("348 位点扩样结果：13 条 strict-pass 候选，形成 8 个家族", fontsize=16, fontweight="bold", color=COLORS["navy"])
    return save_fig(fig, "15_b348_design_update.png")


def create_joint_figure(combo: pd.DataFrame, shortlist: pd.DataFrame) -> Path:
    fig = plt.figure(figsize=(12.6, 7.0))
    gs = fig.add_gridspec(2, 2, hspace=0.35, wspace=0.32)

    ax1 = fig.add_subplot(gs[0, 0])
    tier_counts = combo["joint_priority_tier"].value_counts().reindex(["high", "medium", "low"]).fillna(0)
    ax1.bar(tier_counts.index, tier_counts.values, color=["#2A9D8F", "#E9C46A", "#8C98A4"])
    for i, v in enumerate(tier_counts.values):
        ax1.text(i, v + 2, str(int(v)), ha="center", fontsize=10)
    ax1.set_title("组合优先级分层", fontsize=13)
    ax1.set_ylabel("组合数量")

    ax2 = fig.add_subplot(gs[0, 1])
    palette = {"high": "#2A9D8F", "medium": "#D8A329", "low": "#8C98A4"}
    for tier, sub in combo.groupby("joint_priority_tier"):
        ax2.scatter(sub["peptide_peptide_min_heavy_distance_A"], sub["compatibility_score"], s=38, alpha=0.8, label=tier, color=palette.get(tier, "#8C98A4"))
    ax2.set_title("静态兼容性与肽-肽距离", fontsize=13)
    ax2.set_xlabel("肽-肽最小重原子距离（A）")
    ax2.set_ylabel("兼容性得分")
    ax2.legend(frameon=False, title="层级")

    ax3 = fig.add_subplot(gs[1, 0])
    heat = combo.pivot_table(index="family_348", columns="family_404405", values="compatibility_score", aggfunc="mean")
    sns.heatmap(heat, ax=ax3, cmap="YlGnBu", annot=True, fmt=".3f", cbar_kws={"shrink": 0.7})
    ax3.set_title("家族配对平均得分", fontsize=13)
    ax3.set_xlabel("404/405 家族")
    ax3.set_ylabel("348 家族")

    ax4 = fig.add_subplot(gs[1, 1])
    top = shortlist.head(8).copy().iloc[::-1]
    labels = [f"{row.sequence_404405} + {row.sequence_348}" for row in top.itertuples(index=False)]
    ax4.barh(range(len(top)), top["compatibility_score"], color="#4C78A8")
    ax4.set_yticks(range(len(top)))
    ax4.set_yticklabels([wrap_label(x, 18) for x in labels], fontsize=8)
    ax4.set_title("实验优先 shortlist", fontsize=13)
    ax4.set_xlabel("兼容性得分")

    fig.suptitle("双肽联合预筛：16 x 13 共 208 个组合", fontsize=16, fontweight="bold", color=COLORS["navy"])
    return save_fig(fig, "16_joint_combo_overview.png")


def create_top_pair_structure_figure(pdb_path: Path, seq_404: str, seq_348: str) -> Path:
    parser = PDBParser(QUIET=True)
    structure = parser.get_structure("combo", str(pdb_path))
    chain_coords: list[tuple[str, np.ndarray]] = []
    for chain in structure[0]:
        coords = []
        for residue in chain:
            if "CA" in residue:
                coords.append(residue["CA"].coord)
        if coords:
            chain_coords.append((chain.id, np.asarray(coords, dtype=float)))
    chain_coords.sort(key=lambda item: len(item[1]), reverse=True)
    target_id, target = chain_coords[0]
    pep1_id, pep1 = chain_coords[1]
    pep2_id, pep2 = chain_coords[2]
    center = target.mean(axis=0)
    target = target - center
    pep1 = pep1 - center
    pep2 = pep2 - center

    fig, axes = plt.subplots(1, 2, figsize=(12.4, 5.2))
    views = [("XY 投影", 0, 1), ("XZ 投影", 0, 2)]
    for ax, (title, i, j) in zip(axes, views):
        ax.plot(target[:, i], target[:, j], color="#4C78A8", linewidth=2.0, alpha=0.9, label=f"靶蛋白链 {target_id}")
        ax.plot(pep1[:, i], pep1[:, j], color="#E07A5F", linewidth=3.0, marker="o", markersize=4, label=f"404/405 多肽 {pep1_id}")
        ax.plot(pep2[:, i], pep2[:, j], color="#D8A329", linewidth=3.0, marker="o", markersize=4, label=f"348 多肽 {pep2_id}")
        ax.set_title(title, fontsize=13)
        ax.set_xlabel("X（A）")
        ax.set_ylabel(("Y" if j == 1 else "Z") + "（A）")
        ax.legend(frameon=False, loc="best", fontsize=9)
        ax.set_aspect("equal", adjustable="datalim")
    fig.suptitle(f"最优双肽组合结构：{seq_404} + {seq_348}", fontsize=16, fontweight="bold", color=COLORS["navy"])
    return save_fig(fig, "17_top_pair_structure.png")


def create_shortlist_figure(shortlist: pd.DataFrame) -> Path:
    show = shortlist.head(8).copy()
    fig, ax = plt.subplots(figsize=(12.6, 5.8))
    ax.axis("off")
    columns = ["排名", "404/405 多肽", "348 多肽", "层级", "得分", "距离(A)"]
    table_data = []
    for row in show.itertuples(index=False):
        table_data.append(
            [
                int(row.shortlist_rank),
                row.sequence_404405,
                row.sequence_348,
                row.joint_priority_tier,
                f"{row.compatibility_score:.3f}",
                f"{row.peptide_peptide_min_heavy_distance_A:.2f}",
            ]
        )
    table = ax.table(cellText=table_data, colLabels=columns, loc="center", cellLoc="center", colLoc="center")
    table.auto_set_font_size(False)
    table.set_fontsize(10)
    table.scale(1, 1.8)
    for (r, c), cell in table.get_celld().items():
        if r == 0:
            cell.set_facecolor("#17324D")
            cell.set_text_props(color="white", weight="bold")
        else:
            cell.set_facecolor("#F6F4EF" if r % 2 == 1 else "white")
            if c == 3:
                tier = table_data[r - 1][3]
                if tier == "high":
                    cell.set_text_props(color="#2A9D8F", weight="bold")
                elif tier == "medium":
                    cell.set_text_props(color="#D8A329", weight="bold")
                else:
                    cell.set_text_props(color="#8C98A4", weight="bold")
    ax.set_title("优先验证的双肽组合", fontsize=16, fontweight="bold", color=COLORS["navy"], pad=16)
    return save_fig(fig, "18_joint_shortlist_table.png")


def build_presentation(
    md_stats: dict,
    df404: pd.DataFrame,
    fam404: pd.DataFrame,
    df348: pd.DataFrame,
    fam348: pd.DataFrame,
    combo: pd.DataFrame,
    shortlist: pd.DataFrame,
    assets: dict[str, Path],
) -> None:
    prs = Presentation()
    prs.slide_width = base.PPT_WIDTH
    prs.slide_height = base.PPT_HEIGHT
    blank = prs.slide_layouts[6]

    top_pair = shortlist.iloc[0]

    slide = prs.slides.add_slide(blank)
    base.apply_slide_style(slide, "靶蛋白-病毒蛋白界面解析与双位点多肽设计", "MD -> 热点 -> 设计 -> 组合")
    base.add_textbox(slide, Inches(0.7), Inches(1.0), Inches(7.9), Inches(0.85), "基于 MD 锁定 B404/B405 主热点，并进一步扩展到 B348 次级位点，完成 8-12 aa 多肽设计与双肽组合预筛。", font_size=24, bold=True, color=COLORS["navy"])
    add_metric_box(slide, Inches(0.8), Inches(2.0), Inches(2.3), Inches(1.0), "16", "404/405 strict-pass 候选", COLORS["accent"])
    add_metric_box(slide, Inches(3.35), Inches(2.0), Inches(2.3), Inches(1.0), "13", "348 strict-pass 候选", COLORS["target"])
    add_metric_box(slide, Inches(5.9), Inches(2.0), Inches(2.3), Inches(1.0), "208", "双肽组合数", COLORS["hotspot"])
    add_metric_box(slide, Inches(8.45), Inches(2.0), Inches(2.3), Inches(1.0), "24", "高优先级组合", COLORS["muted_red"])
    base.add_bullets(
        slide,
        Inches(0.9),
        Inches(3.35),
        Inches(7.0),
        Inches(2.7),
        [
            "MD 结果支持先围绕 B404/B405 做主热点设计，再扩展到空间分离的 B348 位点。",
            "Complexa 扩样后获得 404/405 位点 16 条、348 位点 13 条严格过筛候选。",
            "全部 16 x 13 个组合完成联合预筛，当前最优组合为 DKTKTINV + TLRKLAEQLAED。",
        ],
        font_size=18,
    )
    base.add_image_fit(slide, OLD_ASSETS / "01_workflow_overview.png", Inches(8.6), Inches(3.0), width=Inches(4.0))
    add_source_footer(slide, "数据来源：MD 重分析 + Proteina-Complexa 扩样结果 + 双肽联合静态筛选")

    slide = prs.slides.add_slide(blank)
    base.apply_slide_style(slide, "研究逻辑与证据链", "总体框架")
    base.add_image_fit(slide, OLD_ASSETS / "01_workflow_overview.png", Inches(0.6), Inches(1.0), width=Inches(6.0))
    base.add_bullets(
        slide,
        Inches(7.0),
        Inches(1.15),
        Inches(5.5),
        Inches(2.7),
        [
            "第一步：通过 RMSD、界面距离、接触数和氢键占据率重分析 MD。",
            "第二步：将 B404/B405 定位为主热点，并从结构几何上识别 B348 为第二位点。",
            "第三步：围绕两个位点分别进行 8-12 aa 热点引导设计，并用统一标准筛选。",
            "第四步：将 404/405 候选与 348 候选做两两组合，评估是否适合联合使用。",
        ],
        font_size=17,
    )
    base.add_textbox(slide, Inches(7.1), Inches(4.65), Inches(5.0), Inches(1.2), "本版已不再停留于早期 pilot 结果，而是基于扩样后的 404/405、348 全量候选和 208 个组合做出结论。", font_size=15, color=COLORS["ink"])
    add_source_footer(slide, "逻辑主线：热点识别 -> 单位点设计 -> 双位点联合策略")

    slide = prs.slides.add_slide(blank)
    base.apply_slide_style(slide, "MD 稳定性重分析", "0-50 ns 可复现实验轨迹")
    base.add_image_fit(slide, OLD_ASSETS / "02_rmsd_overview.png", Inches(0.55), Inches(0.95), width=Inches(8.3))
    base.add_textbox(slide, Inches(9.0), Inches(1.15), Inches(3.3), Inches(0.35), "核心读数", font_size=18, bold=True, color=COLORS["navy"])
    base.add_bullets(
        slide,
        Inches(9.0),
        Inches(1.55),
        Inches(3.4),
        Inches(2.3),
        [
            f"复合物 RMSD（15-50 ns）：{md_stats['rmsd_complex']['mean']:.2f} +/- {md_stats['rmsd_complex']['std']:.2f} nm",
            f"靶蛋白链 RMSD：{md_stats['rmsd_a']['mean']:.2f} +/- {md_stats['rmsd_a']['std']:.2f} nm",
            f"病毒蛋白链 RMSD：{md_stats['rmsd_b']['mean']:.2f} +/- {md_stats['rmsd_b']['std']:.2f} nm",
            "解释：虽然病毒蛋白链整体波动较大，但界面没有解离，仍可支持后续热点分析。",
        ],
        font_size=15,
    )
    base.add_image_fit(slide, OLD_ASSETS / "04_rg_and_hbond_lifetime.png", Inches(8.8), Inches(4.0), width=Inches(3.8))
    add_source_footer(slide, "分析窗口：15-50 ns；该页用于说明体系可用于后续界面与热点统计")

    slide = prs.slides.add_slide(blank)
    base.apply_slide_style(slide, "界面热点证据", "接触与氢键持续存在")
    base.add_image_fit(slide, OLD_ASSETS / "03_interface_persistence.png", Inches(0.55), Inches(1.0), width=Inches(6.1))
    base.add_image_fit(slide, OLD_ASSETS / "06_hbond_occupancy_hotspots.png", Inches(6.95), Inches(1.0), width=Inches(5.8))
    base.add_bullets(
        slide,
        Inches(0.8),
        Inches(5.6),
        Inches(12.0),
        Inches(1.0),
        [
            f"界面最小距离为 {md_stats['mindist']['mean']:.3f} +/- {md_stats['mindist']['std']:.3f} nm，15-50 ns 内接触始终维持较高水平。",
            f"氢键数为 {md_stats['hbonds_mean']:.2f} +/- {md_stats['hbonds_std']:.2f}；占据率最高的残基集中在 B{md_stats['hotspots'][0]} 和 B{md_stats['hotspots'][1]}。",
        ],
        font_size=16,
    )
    add_source_footer(slide, "结论：B404/B405 是 MD 直接支持的主热点，而不是仅凭静态结构主观选点")

    slide = prs.slides.add_slide(blank)
    base.apply_slide_style(slide, "结构视角下的热点定位", "空间关系解释后续设计")
    base.add_image_fit(slide, OLD_ASSETS / "07_md_structure_overview.png", Inches(0.55), Inches(1.0), width=Inches(5.9))
    base.add_image_fit(slide, OLD_ASSETS / "08_snapshot_strip.png", Inches(6.7), Inches(1.0), width=Inches(5.9))
    base.add_image_fit(slide, OLD_ASSETS / "09_hotspot_network.png", Inches(3.7), Inches(4.35), width=Inches(5.7))
    base.add_textbox(slide, Inches(0.8), Inches(6.55), Inches(12.0), Inches(0.35), "要点：B404/B405 是紧凑热点面，而 B348 与其空间分离，更适合作为第二独立表位进入双肽联合策略。", font_size=15, bold=True, color=COLORS["navy"])
    add_source_footer(slide, "结构解释：404/405 与 348 不是同一个口袋，联合占位在几何上有可行性")

    slide = prs.slides.add_slide(blank)
    base.apply_slide_style(slide, "404/405 扩样设计结果", "16 条 strict-pass 候选")
    base.add_image_fit(slide, assets["design404"], Inches(0.45), Inches(0.95), width=Inches(8.4))
    base.add_textbox(slide, Inches(9.0), Inches(1.1), Inches(3.1), Inches(0.35), "结论", font_size=18, bold=True, color=COLORS["navy"])
    base.add_bullets(
        slide,
        Inches(9.0),
        Inches(1.55),
        Inches(3.3),
        Inches(2.5),
        [
            "扩样后由早期 1 条 pilot 候选扩大到 16 条严格过筛候选。",
            "主家族为 SMEK/MEK 酸性 motif，但 DKTKTINV 代表了另一条高质量独立解。",
            "前列候选普遍同时具备较高肽链置信度和较好的界面置信度。",
        ],
        font_size=15,
    )
    base.add_textbox(slide, Inches(9.0), Inches(4.75), Inches(3.0), Inches(1.3), "优先代表序列：\nDKTKTINV\nSSMEKLEDLINR\nSMEKLDDLIDR\nMEKLEDLIKR", font_size=16, color=COLORS["ink"])
    add_source_footer(slide, "筛选逻辑：先保留所有 strict-pass，再按家族与代表序列组织结果")

    slide = prs.slides.add_slide(blank)
    base.apply_slide_style(slide, "348 扩样设计结果", "13 条 strict-pass 候选")
    base.add_image_fit(slide, assets["design348"], Inches(0.45), Inches(0.95), width=Inches(8.4))
    base.add_textbox(slide, Inches(9.0), Inches(1.1), Inches(3.1), Inches(0.35), "为什么 348 重要", font_size=18, bold=True, color=COLORS["navy"])
    base.add_bullets(
        slide,
        Inches(9.0),
        Inches(1.55),
        Inches(3.3),
        Inches(2.8),
        [
            "348 不是对 404/405 的重复验证，而是第二个独立表位。",
            "三轮扩样后获得 13 条严格过筛候选，且形成 8 个家族，说明不是偶然单点命中。",
            "在联合分析中，TLRKLAEQLAED、TNTKTYTT、NVEIIPVN 反复出现在前列组合中，说明 348 对双肽策略有真实贡献。",
        ],
        font_size=15,
    )
    base.add_textbox(slide, Inches(9.0), Inches(4.85), Inches(3.0), Inches(1.2), "优先代表序列：\nTLRKLAEQLAED\nTNTKTYTT\nNVEIIPVN\nKTVTLEDLLA", font_size=15, color=COLORS["ink"])
    add_source_footer(slide, "核心逻辑：348 使项目从“单热点肽”推进到“两个空间分离位点的联合占位”")

    slide = prs.slides.add_slide(blank)
    base.apply_slide_style(slide, "筛选指标为什么能说明结果更好", "指标解释")
    base.add_textbox(slide, Inches(0.75), Inches(1.0), Inches(5.7), Inches(0.4), "结构/界面置信度指标", font_size=18, bold=True, color=COLORS["navy"])
    base.add_bullets(
        slide,
        Inches(0.8),
        Inches(1.45),
        Inches(5.8),
        Inches(3.6),
        [
            "Binder pLDDT：多肽自身局部结构置信度，越高越好。",
            "iPTM：靶蛋白与多肽界面整体排布置信度，越高越好。",
            "min iPAE / 平均界面 PAE：界面对齐误差，越低越好。",
            "这些指标不能直接等价于亲和力，但能说明“结构上像不像一个可信的结合姿势”。",
        ],
        font_size=16,
    )
    base.add_textbox(slide, Inches(6.75), Inches(1.0), Inches(5.7), Inches(0.4), "热点命中与组合指标", font_size=18, bold=True, color=COLORS["navy"])
    base.add_bullets(
        slide,
        Inches(6.8),
        Inches(1.45),
        Inches(5.7),
        Inches(3.6),
        [
            "hotspot coverage / 最小热点距离 / 接触残基数：判断是否真的打在目标位点上。",
            "compatibility score：双肽组合内部排序分，只在同一批组合里比较。",
            "肽-肽最小距离、hard clash、soft clash：判断两条肽能否同时放到靶蛋白表面。",
            "因此我们看的不是单一高分，而是“位点命中 + 结构可信 + 组合兼容”三者同时成立。",
        ],
        font_size=16,
    )
    add_source_footer(slide, "指标口径：AlphaFold-Multimer 常用 pLDDT / iPTM / PAE；本工作另加入热点命中与组合几何指标")

    slide = prs.slides.add_slide(blank)
    base.apply_slide_style(slide, "官方与同类研究如何使用这些指标", "方法依据")
    box1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.2), Inches(3.85), Inches(4.9))
    box1.fill.solid(); box1.fill.fore_color.rgb = rgb_color("#FFFFFF"); box1.line.color.rgb = rgb_color(COLORS["accent"])
    box2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(4.75), Inches(1.2), Inches(3.85), Inches(4.9))
    box2.fill.solid(); box2.fill.fore_color.rgb = rgb_color("#FFFFFF"); box2.line.color.rgb = rgb_color(COLORS["target"])
    box3 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(1.2), Inches(3.85), Inches(4.9))
    box3.fill.solid(); box3.fill.fore_color.rgb = rgb_color("#FFFFFF"); box3.line.color.rgb = rgb_color(COLORS["hotspot"])
    base.add_textbox(slide, Inches(0.95), Inches(1.45), Inches(3.3), Inches(0.3), "Proteina-Complexa 官方", font_size=18, bold=True, color=COLORS["navy"])
    base.add_bullets(slide, Inches(0.95), Inches(1.85), Inches(3.2), Inches(3.8), [
        "官方 README 明确把 reward model、success filtering、diversity computation 作为流程核心。",
        "官方模型卡强调：模型输出后仍需结合具体任务再测试与验证。",
        "因此本工作采用“批量生成 -> 过滤 -> 分析 -> 按家族整理”的汇报逻辑。"
    ], font_size=14)
    base.add_textbox(slide, Inches(5.0), Inches(1.45), Inches(3.2), Inches(0.3), "AlphaFold 类指标的常见含义", font_size=18, bold=True, color=COLORS["navy"])
    base.add_bullets(slide, Inches(5.0), Inches(1.85), Inches(3.2), Inches(3.8), [
        "pLDDT 常用于判断局部结构是否可信。",
        "iPTM 常用于判断界面排布是否可信。",
        "PAE / iPAE 常用于判断跨链相对定位误差，越低越好。"
    ], font_size=14)
    base.add_textbox(slide, Inches(9.05), Inches(1.45), Inches(3.2), Inches(0.3), "同类 binder 设计研究的做法", font_size=18, bold=True, color=COLORS["navy"])
    base.add_bullets(slide, Inches(9.05), Inches(1.85), Inches(3.2), Inches(3.8), [
        "BindCraft 官方建议先指定热点、跑大量轨迹，再保留所有通过过滤的设计。",
        "官方还建议生成足够多的 pass designs，再从前列少量候选进入实验。",
        "BindCraft 明确指出 iPTM 更像“是否可能结合”的指标，而不是亲和力本身。"
    ], font_size=14)
    add_source_footer(slide, "依据：Proteina-Complexa README / model card；BindCraft README；AlphaFold-Multimer 常用指标解释")

    slide = prs.slides.add_slide(blank)
    base.apply_slide_style(slide, "双肽联合使用预筛", "208 个组合全部完成评分")
    base.add_image_fit(slide, assets["joint"], Inches(0.42), Inches(0.95), width=Inches(8.7))
    base.add_bullets(
        slide,
        Inches(9.3),
        Inches(1.2),
        Inches(3.0),
        Inches(2.6),
        [
            "16 条 404/405 候选与 13 条 348 候选全部完成两两组合。",
            "208 个组合都能通过基础几何检查，但综合单肽质量后只有 24 个属于高优先级。",
            "348 侧最强的组合伙伴是 TLRKLAEQLAED，这一趋势在前列组合中反复出现。",
        ],
        font_size=15,
    )
    base.add_textbox(slide, Inches(9.25), Inches(4.55), Inches(3.0), Inches(1.1), "前 5 个组合都把一个强 404/405 候选与 TLRKLAEQLAED 配对，提示它可能是最稳定的第二位点搭档。", font_size=15, color=COLORS["navy"])
    add_source_footer(slide, "判断标准：几何兼容 + 单肽本身质量 + 位点保留，不直接等同于最终亲和力")

    slide = prs.slides.add_slide(blank)
    base.apply_slide_style(slide, "最佳双肽组合", "当前最优 pair")
    base.add_image_fit(slide, assets["top_pair"], Inches(0.55), Inches(1.0), width=Inches(7.2))
    add_metric_box(slide, Inches(8.4), Inches(1.15), Inches(3.7), Inches(0.95), "DKTKTINV + TLRKLAEQLAED", "最优双肽序列", COLORS["navy"])
    base.add_bullets(
        slide,
        Inches(8.55),
        Inches(2.35),
        Inches(3.3),
        Inches(3.4),
        [
            f"组合得分：{top_pair['compatibility_score']:.3f}",
            f"优先级：{top_pair['joint_priority_tier']}",
            f"404/405 多肽：pLDDT {top_pair['binder_plddt_404405']:.3f}，iPTM {top_pair['i_ptm_404405']:.3f}",
            f"348 多肽：pLDDT {top_pair['binder_plddt_348']:.3f}，iPTM {top_pair['i_ptm_348']:.3f}",
            f"肽-肽最小距离：{top_pair['peptide_peptide_min_heavy_distance_A']:.2f} A；碰撞：0 hard / 0 soft",
            f"跨热点干扰距离：{top_pair['cross_404405_peptide_to_348_hotspot_min_A']:.2f} A 和 {top_pair['cross_348_peptide_to_404405_hotspot_min_A']:.2f} A",
        ],
        font_size=14,
    )
    add_source_footer(slide, "当前解读：这是“最适合进入下一轮验证”的组合，不等同于已证明亲和力最高")

    slide = prs.slides.add_slide(blank)
    base.apply_slide_style(slide, "后续实验优先列表", "适合优先验证的组合")
    base.add_image_fit(slide, assets["shortlist"], Inches(0.55), Inches(1.05), width=Inches(8.0))
    base.add_bullets(
        slide,
        Inches(8.9),
        Inches(1.3),
        Inches(3.1),
        Inches(2.6),
        [
            "第一优先组：DKTKTINV 分别配对 TLRKLAEQLAED、TNTKTYTT、NVEIIPVN。",
            "第二优先组：SMEK 家族代表序列与 TLRKLAEQLAED 配对。",
            "这组组合最适合进入下一步局部精修、短程约束 MD 或实验前复核。",
        ],
        font_size=15,
    )
    base.add_textbox(slide, Inches(8.95), Inches(4.95), Inches(3.0), Inches(1.2), "建议下一步：\n1. 快速局部最小化\n2. 短程约束 MD\n3. 结合能或接触持续性复核", font_size=15, color=COLORS["ink"])
    add_source_footer(slide, "shortlist 用途：缩小验证范围，不是提前替代实验结论")

    slide = prs.slides.add_slide(blank)
    base.apply_slide_style(slide, "结论与建议", "更新后的核心信息")
    base.add_bullets(
        slide,
        Inches(0.85),
        Inches(1.2),
        Inches(7.3),
        Inches(4.7),
        [
            "MD 重分析支持 B404/B405 为主热点，B348 为第二个空间分离位点。",
            "在严格 8-12 aa 条件下，Complexa 扩样后获得 404/405 位点 16 条、348 位点 13 条严格过筛候选。",
            "348 的意义不在于重复命中，而在于它提供了可与主热点并行占位的第二抓手。",
            "208 个双肽组合中有 24 个高优先级组合，当前最优为 DKTKTINV + TLRKLAEQLAED。",
            "下一步最合理的是围绕 shortlist 做精修与小规模验证，而不是再回到盲目扩样。",
        ],
        font_size=18,
    )
    summary_shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.4), Inches(1.4), Inches(3.8), Inches(3.3))
    summary_shape.fill.solid()
    summary_shape.fill.fore_color.rgb = rgb_color("#FFFFFF")
    summary_shape.line.color.rgb = rgb_color(COLORS["hotspot"])
    base.add_textbox(slide, Inches(8.7), Inches(1.75), Inches(3.2), Inches(0.35), "一句话总结", font_size=18, bold=True, color=COLORS["navy"])
    base.add_textbox(slide, Inches(8.7), Inches(2.25), Inches(3.0), Inches(1.8), "本课题已经从“找一个短肽”推进到“围绕两个独立位点建立双肽联合策略”，并筛出了可直接进入下一轮验证的组合。", font_size=16, color=COLORS["ink"])
    add_source_footer(slide, "汇总结论基于：MD 热点证据 + 单位点 strict-pass + 双肽组合静态预筛")

    prs.save(str(PPTX_PATH))


def main():
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    md_stats = compute_md_stats()
    df404, fam404, df348, fam348, combo, shortlist = read_tables()
    assets = {
        "design404": create_404405_figure(df404, fam404),
        "design348": create_348_figure(df348, fam348),
        "joint": create_joint_figure(combo, shortlist),
        "top_pair": create_top_pair_structure_figure(Path(shortlist.iloc[0]["merged_combo_pdb"]), shortlist.iloc[0]["sequence_404405"], shortlist.iloc[0]["sequence_348"]),
        "shortlist": create_shortlist_figure(shortlist),
    }
    build_presentation(md_stats, df404, fam404, df348, fam348, combo, shortlist, assets)

    summary_lines = [
        "LY 课题汇报更新版摘要",
        f"上一版 PPT：{OLD_PPT}",
        f"新版 PPT：{PPTX_PATH}",
        f"素材目录：{ASSET_DIR}",
        f"MD 分析窗口：{base.ANALYSIS_START_NS:.0f}-{base.ANALYSIS_END_NS:.0f} ns",
        f"MD 支持的主热点：B{md_stats['hotspots'][0]}、B{md_stats['hotspots'][1]}",
        f"404/405 strict-pass 候选数：{len(df404)}",
        f"348 strict-pass 候选数：{len(df348)}",
        f"双肽组合总数：{len(combo)}",
        f"高优先级组合数：{int((combo['joint_priority_tier'] == 'high').sum())}",
        f"当前最优组合：{shortlist.iloc[0]['sequence_404405']} + {shortlist.iloc[0]['sequence_348']}",
        f"当前最优组合得分：{float(shortlist.iloc[0]['compatibility_score']):.4f}",
    ]
    SUMMARY_PATH.write_text("\n".join(summary_lines), encoding="utf-8")


if __name__ == "__main__":
    main()
