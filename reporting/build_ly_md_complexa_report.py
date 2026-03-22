from __future__ import annotations

import csv
import re
from collections import Counter, defaultdict
from dataclasses import dataclass
from pathlib import Path

import matplotlib
import matplotlib.font_manager as fm

matplotlib.use("Agg")

import matplotlib.gridspec as gridspec
import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import seaborn as sns
from Bio.PDB import PDBParser
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"D:\Project\cw\ly")
MD_DIR = ROOT / "fold_pig_vapa_n" / "MD" / "newMD"
PLOT_DIR = MD_DIR / "plot"
PDB_FRAME_DIR = MD_DIR / "pdb"
DESIGN_DIR = ROOT / "Complexa_B404_B405_secondary_screening_bulk512"
DESIGN_FINAL_PDB = DESIGN_DIR / "final_complexes" / "01_job_0_n_167_id_27_single_orig7.pdb"
ASSET_DIR = ROOT / "LY_MD_Complexa_Report_assets"
PPTX_PATH = ROOT / "LY_MD_Complexa_Report_20260321.pptx"
SUMMARY_PATH = ROOT / "LY_MD_Complexa_Report_20260321_summary.txt"
PPT_WIDTH = Inches(13.333)
PPT_HEIGHT = Inches(7.5)

CHAIN_A_LEN = 118
HOTSPOT_MD_RAW = [522, 523]
HOTSPOT_MD_LABELS = ["B404", "B405"]
HOTSPOT_DESIGN = [140, 141]
HOTSPOT_DESIGN_LABELS = ["B404", "B405"]
ANALYSIS_START_NS = 15.0
ANALYSIS_END_NS = 50.0

COLORS = {
    "navy": "#17324D",
    "ink": "#22303C",
    "slate": "#51606D",
    "grid": "#D8DEE6",
    "bg": "#F6F4EF",
    "target": "#4C78A8",
    "virus": "#E07A5F",
    "complex": "#2F4858",
    "hotspot": "#D8A329",
    "accent": "#2A9D8F",
    "muted_red": "#C8553D",
    "gold": "#E9C46A",
    "gray": "#8C98A4",
}

FONT_FAMILY = "Microsoft YaHei"

plt.rcParams.update(
    {
        "figure.facecolor": "white",
        "axes.facecolor": "white",
        "axes.edgecolor": COLORS["slate"],
        "axes.labelcolor": COLORS["ink"],
        "xtick.color": COLORS["ink"],
        "ytick.color": COLORS["ink"],
        "grid.color": COLORS["grid"],
        "font.family": [FONT_FAMILY],
        "font.sans-serif": [FONT_FAMILY, "SimHei", "DejaVu Sans"],
        "axes.titleweight": "bold",
        "axes.titlecolor": COLORS["navy"],
        "axes.unicode_minus": False,
    }
)
sns.set_style("whitegrid")
for font_path in [r"C:\Windows\Fonts\msyh.ttc", r"C:\Windows\Fonts\msyhbd.ttc", r"C:\Windows\Fonts\simhei.ttf"]:
    try:
        fm.fontManager.addfont(font_path)
    except Exception:
        pass
plt.rcParams["font.family"] = FONT_FAMILY
plt.rcParams["font.sans-serif"] = [FONT_FAMILY, "SimHei", "DejaVu Sans"]
plt.rcParams["axes.unicode_minus"] = False

AA3_TO_1 = {
    "ALA": "A",
    "ARG": "R",
    "ASN": "N",
    "ASP": "D",
    "CYS": "C",
    "GLN": "Q",
    "GLU": "E",
    "GLY": "G",
    "HIS": "H",
    "ILE": "I",
    "LEU": "L",
    "LYS": "K",
    "MET": "M",
    "PHE": "F",
    "PRO": "P",
    "SER": "S",
    "THR": "T",
    "TRP": "W",
    "TYR": "Y",
    "VAL": "V",
}


@dataclass
class XvgSeries:
    time_ns: np.ndarray
    value: np.ndarray


@dataclass
class HbondRecord:
    occupancy: float
    donor_desc: str
    mapping_desc: str
    left_chain: str
    left_resid: int
    right_chain: str
    right_resid: int


def ensure_dirs() -> None:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)


def read_xvg(path: Path, scale_time_ps_to_ns: bool = True) -> XvgSeries:
    times: list[float] = []
    values: list[float] = []
    with path.open() as handle:
        for line in handle:
            if not line or line[0] in "#@":
                continue
            parts = line.split()
            if len(parts) < 2:
                continue
            time_val = float(parts[0])
            if scale_time_ps_to_ns:
                time_val /= 1000.0
            times.append(time_val)
            values.append(float(parts[1]))
    return XvgSeries(np.asarray(times), np.asarray(values))


def read_xvg_two_series(path: Path, scale_time_ps_to_ns: bool = True) -> tuple[np.ndarray, np.ndarray, np.ndarray]:
    times: list[float] = []
    v1: list[float] = []
    v2: list[float] = []
    with path.open() as handle:
        for line in handle:
            if not line or line[0] in "#@":
                continue
            parts = line.split()
            if len(parts) < 3:
                continue
            time_val = float(parts[0])
            if scale_time_ps_to_ns:
                time_val /= 1000.0
            times.append(time_val)
            v1.append(float(parts[1]))
            v2.append(float(parts[2]))
    return np.asarray(times), np.asarray(v1), np.asarray(v2)


def moving_average(values: np.ndarray, window: int = 101) -> np.ndarray:
    if window <= 1 or len(values) < window:
        return values.copy()
    if window % 2 == 0:
        window += 1
    kernel = np.ones(window) / window
    return np.convolve(values, kernel, mode="same")


def compute_stats(series: XvgSeries, start_ns: float, end_ns: float | None = None) -> dict[str, float]:
    mask = series.time_ns >= start_ns
    if end_ns is not None:
        mask &= series.time_ns <= end_ns
    values = series.value[mask]
    times = series.time_ns[mask]
    if len(values) == 0:
        values = series.value
        times = series.time_ns
    slope = 0.0 if len(values) < 2 else float(np.polyfit(times, values, deg=1)[0])
    return {
        "mean": float(np.mean(values)),
        "std": float(np.std(values)),
        "min": float(np.min(values)),
        "max": float(np.max(values)),
        "slope": slope,
        "n": int(len(values)),
    }


def parse_hbond_records(path: Path) -> list[HbondRecord]:
    records: list[HbondRecord] = []
    pattern = re.compile(r"#\s+([0-9.]+)%\s+(.+?)\s+\((.+)\)")
    chain_res_pattern = re.compile(r"([AB])(\d+)")
    for line in path.read_text(encoding="utf-8", errors="ignore").splitlines():
        match = pattern.match(line.strip())
        if not match:
            continue
        occupancy = float(match.group(1))
        donor_desc = match.group(2).strip()
        mapping = match.group(3).strip()
        left, right = [token.strip() for token in mapping.split("->")]
        left_match = chain_res_pattern.match(left)
        right_match = chain_res_pattern.match(right)
        if not left_match or not right_match:
            continue
        records.append(
            HbondRecord(
                occupancy=occupancy,
                donor_desc=donor_desc,
                mapping_desc=mapping,
                left_chain=left_match.group(1),
                left_resid=int(left_match.group(2)),
                right_chain=right_match.group(1),
                right_resid=int(right_match.group(2)),
            )
        )
    return records


def summarize_hbond_records(records: list[HbondRecord]) -> tuple[pd.DataFrame, pd.DataFrame]:
    atom_df = pd.DataFrame(
        [{"occupancy": r.occupancy, "interaction": r.mapping_desc, "description": r.donor_desc} for r in records]
    ).sort_values("occupancy", ascending=False)
    residue_max = defaultdict(float)
    residue_sum = defaultdict(float)
    for record in records:
        for chain, resid in ((record.left_chain, record.left_resid), (record.right_chain, record.right_resid)):
            if chain == "B":
                residue_sum[resid] += record.occupancy
                residue_max[resid] = max(residue_max[resid], record.occupancy)
    residue_df = (
        pd.DataFrame(
            [{"residue": resid, "max_occupancy": occ, "sum_occupancy": residue_sum[resid]} for resid, occ in residue_max.items()]
        )
        .sort_values("max_occupancy", ascending=False)
        .reset_index(drop=True)
    )
    return atom_df, residue_df


def load_design_table() -> pd.DataFrame:
    return pd.read_csv(DESIGN_DIR / "all_ranked_with_hotspot_metrics.csv")


def compute_design_summary(df: pd.DataFrame) -> dict[str, object]:
    hotspot_covered = int((df["hotspot_coverage"] == 2).sum())
    binder_confident = int(((df["hotspot_coverage"] == 2) & (df["binder_plddt"] >= 0.58)).sum())
    iptm_good = int(((df["hotspot_coverage"] == 2) & (df["binder_plddt"] >= 0.58) & (df["i_ptm"] >= 0.20)).sum())
    final_pass = int((df["secondary_pass"] == "yes").sum())
    length_counts = Counter(df["binder_length"])
    return {
        "generated": int(len(df)),
        "hotspot_covered": hotspot_covered,
        "binder_confident": binder_confident,
        "iptm_good": iptm_good,
        "final_pass": final_pass,
        "length_counts": dict(sorted(length_counts.items())),
    }


def split_md_residue(raw_resid: int) -> tuple[str, int]:
    if raw_resid <= CHAIN_A_LEN:
        return "A", raw_resid
    return "B", raw_resid - CHAIN_A_LEN


def load_md_projection_data(pdb_path: Path) -> dict[str, list[dict[str, object]]]:
    parser = PDBParser(QUIET=True)
    structure = parser.get_structure("md", str(pdb_path))
    model = next(structure.get_models())
    chain_data = {"A": [], "B": []}
    has_explicit_ab = {"A", "B"}.issubset({chain.id for chain in model})
    for residue in model.get_residues():
        if residue.id[0] != " " or "CA" not in residue:
            continue
        raw_resid = residue.id[1]
        if has_explicit_ab and residue.get_parent().id in {"A", "B"}:
            chain_label = residue.get_parent().id
            mapped_resid = raw_resid
        else:
            chain_label, mapped_resid = split_md_residue(raw_resid)
        chain_data[chain_label].append(
            {
                "raw_resid": raw_resid,
                "resid": mapped_resid,
                "resname": residue.get_resname(),
                "coord": residue["CA"].coord.astype(float),
            }
        )
    return chain_data


def load_design_projection_data(pdb_path: Path) -> dict[str, list[dict[str, object]]]:
    parser = PDBParser(QUIET=True)
    structure = parser.get_structure("design", str(pdb_path))
    model = next(structure.get_models())
    out: dict[str, list[dict[str, object]]] = {}
    for chain in model:
        records: list[dict[str, object]] = []
        for residue in chain:
            if residue.id[0] != " " or "CA" not in residue:
                continue
            records.append(
                {
                    "resid": residue.id[1],
                    "resname": residue.get_resname(),
                    "coord": residue["CA"].coord.astype(float),
                }
            )
        out[chain.id] = records
    return out


def coords_from_records(records: list[dict[str, object]]) -> np.ndarray:
    return np.asarray([record["coord"] for record in records], dtype=float)


def kabsch_align(mobile: np.ndarray, reference: np.ndarray) -> tuple[np.ndarray, np.ndarray, np.ndarray]:
    mobile_center = mobile.mean(axis=0)
    reference_center = reference.mean(axis=0)
    mobile_shift = mobile - mobile_center
    ref_shift = reference - reference_center
    cov = mobile_shift.T @ ref_shift
    v, _, wt = np.linalg.svd(cov)
    d = np.sign(np.linalg.det(v @ wt))
    correction = np.diag([1.0, 1.0, d])
    rotation = v @ correction @ wt
    aligned = mobile_shift @ rotation + reference_center
    translation = reference_center - mobile_center @ rotation
    return aligned, rotation, translation


def project_points(points: np.ndarray, basis: np.ndarray, center: np.ndarray) -> np.ndarray:
    return (points - center) @ basis.T


def build_projection_basis(points: np.ndarray) -> tuple[np.ndarray, np.ndarray]:
    center = points.mean(axis=0)
    centered = points - center
    _, _, vh = np.linalg.svd(centered, full_matrices=False)
    basis = vh[:2]
    return basis, center


def style_axis(ax: plt.Axes) -> None:
    ax.grid(True, alpha=0.25)
    for spine in ["top", "right"]:
        ax.spines[spine].set_visible(False)


def save_fig(fig: plt.Figure, name: str, dpi: int = 220) -> Path:
    path = ASSET_DIR / name
    fig.savefig(path, dpi=dpi, bbox_inches="tight")
    plt.close(fig)
    return path


def create_workflow_figure() -> Path:
    fig, ax = plt.subplots(figsize=(12.5, 4.4))
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis("off")
    steps = [
        ("1", "复合体系构建", "目标蛋白 + 病毒蛋白\n构建复合物并溶剂化"),
        ("2", "生产模拟", "当前可复现实验数据\n为 50 ns 生产轨迹"),
        ("3", "界面稳定性评估", "RMSD / Rg / 距离 / 接触数\n确认持续界面存在"),
        ("4", "热点识别", "15-50 ns 氢键网络\n锁定 Chain B 404/405"),
        ("5", "多肽设计", "Proteina-Complexa\n8-12 aa 候选筛选"),
    ]
    xs = np.linspace(0.08, 0.92, len(steps))
    width = 0.16
    for idx, (step_id, title, subtitle) in enumerate(steps):
        left = xs[idx] - width / 2
        box = plt.Rectangle((left, 0.28), width, 0.46, facecolor=COLORS["bg"], edgecolor=COLORS["navy"], linewidth=2)
        ax.add_patch(box)
        circ = plt.Circle((xs[idx], 0.8), 0.038, facecolor=COLORS["navy"], edgecolor="white", linewidth=1.5)
        ax.add_patch(circ)
        ax.text(xs[idx], 0.8, step_id, ha="center", va="center", color="white", fontsize=13, fontweight="bold")
        ax.text(xs[idx], 0.64, title, ha="center", va="center", fontsize=13, color=COLORS["navy"], fontweight="bold")
        ax.text(xs[idx], 0.46, subtitle, ha="center", va="center", fontsize=10.5, color=COLORS["ink"], linespacing=1.4)
        if idx < len(steps) - 1:
            ax.annotate("", xy=(xs[idx + 1] - width / 2, 0.51), xytext=(xs[idx] + width / 2, 0.51), arrowprops=dict(
                arrowstyle="->", color=COLORS["muted_red"], lw=2.5
            ))
    ax.text(0.02, 0.95, "研究逻辑总览", fontsize=20, fontweight="bold", color=COLORS["navy"], ha="left")
    ax.text(
        0.02,
        0.08,
        "注：当前 newMD 目录中的生产轨迹长度为 50 ns，因此本次重分析与图件均基于 0-50 ns；氢键分析窗口沿用现有结果的 15-50 ns。",
        fontsize=10.5,
        color=COLORS["slate"],
        ha="left",
    )
    return save_fig(fig, "01_workflow_overview.png")


def create_rmsd_figure(rmsd_complex: XvgSeries, rmsd_a: XvgSeries, rmsd_b: XvgSeries) -> Path:
    fig, axes = plt.subplots(3, 1, figsize=(11.5, 10), sharex=True)
    series = [
        ("整体复合物 RMSD", rmsd_complex, COLORS["complex"], "复合物整体位移较大，但后期波动收敛到较窄区间"),
        ("Chain A RMSD", rmsd_a, COLORS["target"], "Chain A 始终稳定，说明一侧结构框架保持完好"),
        ("Chain B RMSD", rmsd_b, COLORS["virus"], "Chain B 发生较大构象重排，更符合柔性界面蛋白行为"),
    ]
    for ax, (title, data, color, note) in zip(axes, series):
        smooth = moving_average(data.value, 121)
        ax.plot(data.time_ns, data.value, color=color, alpha=0.25, linewidth=1.2)
        ax.plot(data.time_ns, smooth, color=color, linewidth=2.4)
        ax.axvspan(ANALYSIS_START_NS, ANALYSIS_END_NS, color=COLORS["gold"], alpha=0.13)
        stats = compute_stats(data, ANALYSIS_START_NS, ANALYSIS_END_NS)
        ax.text(
            0.015,
            0.93,
            f"{ANALYSIS_START_NS:.0f}-{ANALYSIS_END_NS:.0f} ns: mean {stats['mean']:.3f} nm | SD {stats['std']:.3f} nm\n{note}",
            transform=ax.transAxes,
            va="top",
            fontsize=9.8,
            bbox={"facecolor": "white", "edgecolor": color, "alpha": 0.92, "boxstyle": "round,pad=0.35"},
        )
        ax.set_title(title, fontsize=13.5, pad=10)
        ax.set_ylabel("RMSD (nm)")
        style_axis(ax)
    axes[-1].set_xlabel("Time (ns)")
    fig.suptitle("MD 全程构象演化：整体重排明显，但局部稳定界面持续存在", fontsize=18, fontweight="bold", color=COLORS["navy"])
    return save_fig(fig, "02_rmsd_overview.png")


def create_interface_persistence_figure(mindist: XvgSeries, contacts: XvgSeries, hb_time: np.ndarray, hbonds: np.ndarray) -> Path:
    fig, axes = plt.subplots(3, 1, figsize=(11.5, 10), sharex=True)
    panels = [
        ("最小界面距离", mindist.time_ns, mindist.value, COLORS["navy"], "Distance (nm)", f"后 10 ns 平均 {compute_stats(mindist, 40, 50)['mean']:.3f} nm"),
        ("< 0.6 nm 接触数", contacts.time_ns, contacts.value, COLORS["accent"], "Contacts", f"后 10 ns 平均 {compute_stats(contacts, 40, 50)['mean']:.0f}"),
        ("氢键数（15-50 ns）", hb_time, hbonds, COLORS["muted_red"], "H-bonds", f"15-50 ns 平均 {np.mean(hbonds):.1f}"),
    ]
    for ax, (title, xs, ys, color, ylabel, summary) in zip(axes, panels):
        smooth = moving_average(np.asarray(ys), 121 if len(ys) > 200 else 31)
        ax.plot(xs, ys, color=color, alpha=0.22, linewidth=1.0)
        ax.plot(xs, smooth, color=color, linewidth=2.4)
        ax.axvspan(ANALYSIS_START_NS, ANALYSIS_END_NS, color=COLORS["gold"], alpha=0.13)
        ax.text(
            0.015,
            0.9,
            summary,
            transform=ax.transAxes,
            va="top",
            fontsize=10,
            bbox={"facecolor": "white", "edgecolor": color, "alpha": 0.92, "boxstyle": "round,pad=0.3"},
        )
        ax.set_title(title, fontsize=13)
        ax.set_ylabel(ylabel)
        style_axis(ax)
    axes[-1].set_xlabel("Time (ns)")
    fig.suptitle("界面持续存在：距离、接触数与氢键数量在分析窗口内保持稳定", fontsize=18, fontweight="bold", color=COLORS["navy"])
    return save_fig(fig, "03_interface_persistence.png")


def create_rg_lifetime_figure(rg: XvgSeries, hb_lifetime_time_ps: np.ndarray, hb_lifetime_prob: np.ndarray, hb_lifetime_weighted: np.ndarray) -> Path:
    fig = plt.figure(figsize=(11.8, 5.5))
    gs = gridspec.GridSpec(1, 2, figure=fig, width_ratios=[1.15, 1.0])
    ax1 = fig.add_subplot(gs[0, 0])
    ax2 = fig.add_subplot(gs[0, 1])

    smooth = moving_average(rg.value, 121)
    ax1.plot(rg.time_ns, rg.value, color=COLORS["complex"], alpha=0.2, linewidth=1.0)
    ax1.plot(rg.time_ns, smooth, color=COLORS["complex"], linewidth=2.6)
    ax1.axvspan(ANALYSIS_START_NS, ANALYSIS_END_NS, color=COLORS["gold"], alpha=0.13)
    rg_stats = compute_stats(rg, 40, 50)
    ax1.text(
        0.02,
        0.92,
        f"40-50 ns: mean {rg_stats['mean']:.3f} nm | SD {rg_stats['std']:.3f} nm",
        transform=ax1.transAxes,
        fontsize=10,
        va="top",
        bbox={"facecolor": "white", "edgecolor": COLORS["complex"], "alpha": 0.92, "boxstyle": "round,pad=0.3"},
    )
    ax1.set_title("整体紧致度（Rg）", fontsize=13.5)
    ax1.set_xlabel("Time (ns)")
    ax1.set_ylabel("Radius of gyration (nm)")
    style_axis(ax1)

    ax2.plot(hb_lifetime_time_ps, hb_lifetime_prob, color=COLORS["virus"], linewidth=2.2, label="p(t)")
    ax2.plot(hb_lifetime_time_ps, hb_lifetime_weighted, color=COLORS["accent"], linewidth=2.0, label="t·p(t)")
    ax2.set_title("氢键寿命分布", fontsize=13.5)
    ax2.set_xlabel("Lifetime (ps)")
    ax2.set_ylabel("Probability")
    ax2.set_xlim(0, min(1000, hb_lifetime_time_ps.max()))
    ax2.legend(frameon=False)
    style_axis(ax2)

    fig.suptitle("构象紧致度与氢键寿命支持一个持久但非刚性的界面", fontsize=18, fontweight="bold", color=COLORS["navy"])
    return save_fig(fig, "04_rg_and_hbond_lifetime.png")


def create_rmsf_figure(rmsf_a: XvgSeries, rmsf_b: XvgSeries) -> Path:
    fig, axes = plt.subplots(2, 1, figsize=(11.5, 8.2), sharex=False)
    axes[0].plot(rmsf_a.time_ns, rmsf_a.value, color=COLORS["target"], linewidth=1.8)
    axes[0].fill_between(rmsf_a.time_ns, rmsf_a.value, color=COLORS["target"], alpha=0.12)
    axes[0].set_title("Chain A 残基柔性分布", fontsize=13)
    axes[0].set_ylabel("RMSF (nm)")
    style_axis(axes[0])

    axes[1].plot(rmsf_b.time_ns, rmsf_b.value, color=COLORS["virus"], linewidth=1.8)
    axes[1].fill_between(rmsf_b.time_ns, rmsf_b.value, color=COLORS["virus"], alpha=0.12)
    for resid, label in zip([404, 405], ["B404", "B405"]):
        match_idx = np.where(rmsf_b.time_ns == resid)[0]
        if len(match_idx) == 0:
            continue
        idx = int(match_idx[0])
        val = float(rmsf_b.value[idx])
        axes[1].scatter([resid], [val], color=COLORS["hotspot"], s=70, zorder=4)
        axes[1].text(resid, val + 0.02, label, color=COLORS["hotspot"], fontsize=11, ha="center", fontweight="bold")
    axes[1].set_title("Chain B 残基柔性分布（标出 404/405 热点）", fontsize=13)
    axes[1].set_xlabel("Residue index")
    axes[1].set_ylabel("RMSF (nm)")
    style_axis(axes[1])

    fig.suptitle("柔性分布显示热点位点位于可接近且持续参与结合的表面区域", fontsize=18, fontweight="bold", color=COLORS["navy"])
    return save_fig(fig, "05_rmsf_hotspots.png")


def create_hbond_occupancy_figure(atom_df: pd.DataFrame, residue_df: pd.DataFrame) -> Path:
    fig = plt.figure(figsize=(12, 6.4))
    gs = gridspec.GridSpec(1, 2, figure=fig, width_ratios=[1.25, 0.9])
    ax1 = fig.add_subplot(gs[0, 0])
    ax2 = fig.add_subplot(gs[0, 1])

    top_atom = atom_df.head(10).iloc[::-1]
    colors = [COLORS["hotspot"] if "B404" in inter or "B405" in inter else COLORS["virus"] for inter in top_atom["interaction"]]
    ax1.barh(top_atom["interaction"], top_atom["occupancy"], color=colors)
    ax1.set_xlabel("Occupancy (%)")
    ax1.set_title("稳定界面氢键（Top 10）", fontsize=13.5)
    style_axis(ax1)

    top_res = residue_df.head(8).iloc[::-1]
    bar_colors = [COLORS["hotspot"] if resid in (404, 405) else COLORS["accent"] for resid in top_res["residue"]]
    ax2.barh([f"B{int(x)}" for x in top_res["residue"]], top_res["max_occupancy"], color=bar_colors)
    ax2.set_xlabel("Max occupancy (%)")
    ax2.set_title("目标蛋白热点残基占据率", fontsize=13.5)
    style_axis(ax2)
    ax2.text(
        0.02,
        0.06,
        "B404/B405 在当前结果中分别对应最高和第二高\n稳定氢键占据率，支持其作为设计热点。",
        transform=ax2.transAxes,
        fontsize=10,
        color=COLORS["ink"],
        bbox={"facecolor": "white", "edgecolor": COLORS["hotspot"], "alpha": 0.9, "boxstyle": "round,pad=0.3"},
    )

    fig.suptitle("氢键网络重分析：404/405 位点是界面最稳定的相互作用中心", fontsize=18, fontweight="bold", color=COLORS["navy"])
    return save_fig(fig, "06_hbond_occupancy_hotspots.png")


def plot_chain_trace(ax: plt.Axes, coords: np.ndarray, color: str, linewidth: float = 2.4) -> None:
    ax.plot(coords[:, 0], coords[:, 1], color=color, linewidth=linewidth, solid_capstyle="round")


def create_md_structure_figure() -> Path:
    reference = load_md_projection_data(PLOT_DIR / "ref_structure_matched.pdb")
    coords_a = coords_from_records(reference["A"])
    coords_b = coords_from_records(reference["B"])
    basis, center = build_projection_basis(np.vstack([coords_a, coords_b]))
    proj_a = project_points(coords_a, basis, center)
    proj_b = project_points(coords_b, basis, center)

    fig, ax = plt.subplots(figsize=(8.5, 7.5))
    plot_chain_trace(ax, proj_a, COLORS["target"], 3.0)
    plot_chain_trace(ax, proj_b, COLORS["virus"], 3.0)
    for raw_resid, label in zip(HOTSPOT_MD_RAW, HOTSPOT_MD_LABELS):
        _, mapped_resid = split_md_residue(raw_resid)
        idx = next(i for i, record in enumerate(reference["B"]) if record["resid"] == mapped_resid)
        xy = proj_b[idx]
        ax.scatter([xy[0]], [xy[1]], s=180, color=COLORS["hotspot"], edgecolors="white", linewidths=1.4, zorder=5)
        ax.text(xy[0] + 0.15, xy[1] + 0.15, label, fontsize=11.5, color=COLORS["navy"], fontweight="bold")
    ax.text(proj_a[:, 0].min(), proj_a[:, 1].max() + 1.0, "Chain A", color=COLORS["target"], fontsize=12, fontweight="bold")
    ax.text(proj_b[:, 0].max() - 2.0, proj_b[:, 1].min() - 0.8, "Chain B", color=COLORS["virus"], fontsize=12, fontweight="bold")
    ax.set_title("参考构象中的复合物整体形态与 404/405 热点位置", fontsize=16, fontweight="bold", color=COLORS["navy"])
    ax.set_xticks([])
    ax.set_yticks([])
    for spine in ax.spines.values():
        spine.set_visible(False)
    return save_fig(fig, "07_md_structure_overview.png")


def create_snapshot_strip_figure() -> Path:
    ref_data = load_md_projection_data(PLOT_DIR / "ref_structure_matched.pdb")
    ref_a = coords_from_records(ref_data["A"])
    ref_b = coords_from_records(ref_data["B"])
    basis, center = build_projection_basis(np.vstack([ref_a, ref_b]))
    frames = [(0, "15 ns"), (1000, "25 ns"), (2000, "35 ns"), (3500, "50 ns")]
    fig, axes = plt.subplots(1, 4, figsize=(15.6, 4.8), sharex=True, sharey=True)
    for ax, (frame_idx, label) in zip(axes, frames):
        frame_data = load_md_projection_data(PDB_FRAME_DIR / f"frame{frame_idx}.pdb")
        mov_a = coords_from_records(frame_data["A"])
        mov_b = coords_from_records(frame_data["B"])
        _, rotation, translation = kabsch_align(mov_a, ref_a)
        aligned_a = (mov_a @ rotation) + translation
        aligned_b = (mov_b @ rotation) + translation
        proj_a = project_points(aligned_a, basis, center)
        proj_b = project_points(aligned_b, basis, center)
        plot_chain_trace(ax, proj_a, COLORS["target"], 2.6)
        plot_chain_trace(ax, proj_b, COLORS["virus"], 2.6)
        for raw_resid in HOTSPOT_MD_RAW:
            _, mapped_resid = split_md_residue(raw_resid)
            idx = next(i for i, record in enumerate(frame_data["B"]) if record["resid"] == mapped_resid)
            xy = proj_b[idx]
            ax.scatter([xy[0]], [xy[1]], s=55, color=COLORS["hotspot"], edgecolors="white", linewidths=0.8, zorder=5)
        ax.set_title(label, fontsize=12.5, color=COLORS["navy"], fontweight="bold")
        ax.set_xticks([])
        ax.set_yticks([])
        for spine in ax.spines.values():
            spine.set_visible(False)
    fig.suptitle("15-50 ns 结构快照：整体重排明显，但热点区域持续暴露并保持可结合状态", fontsize=18, fontweight="bold", color=COLORS["navy"])
    return save_fig(fig, "08_snapshot_strip.png")


def create_hotspot_network_figure(records: list[HbondRecord]) -> Path:
    hotspot_links = []
    for record in records:
        if {record.left_chain, record.right_chain} != {"A", "B"}:
            continue
        if record.left_chain == "B":
            target_res, partner_res = record.left_resid, record.right_resid
        else:
            target_res, partner_res = record.right_resid, record.left_resid
        if target_res in {404, 405, 348, 406, 400, 418}:
            hotspot_links.append((target_res, partner_res, record.occupancy))
    target_levels = sorted({item[0] for item in hotspot_links}, key=lambda x: (-max(v for t, _, v in hotspot_links if t == x), x))
    partner_levels = sorted({item[1] for item in hotspot_links}, key=lambda x: (-max(v for _, p, v in hotspot_links if p == x), x))
    target_y = {res: 0.85 - idx * 0.14 for idx, res in enumerate(target_levels)}
    partner_y = {res: 0.85 - idx * 0.12 for idx, res in enumerate(partner_levels)}

    fig, ax = plt.subplots(figsize=(10.5, 6.6))
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.axis("off")
    ax.text(0.17, 0.95, "目标蛋白热点残基（Chain B）", fontsize=14, fontweight="bold", color=COLORS["virus"], ha="center")
    ax.text(0.83, 0.95, "对侧稳定配对残基（Chain A）", fontsize=14, fontweight="bold", color=COLORS["target"], ha="center")
    for res, y in target_y.items():
        color = COLORS["hotspot"] if res in {404, 405} else COLORS["virus"]
        ax.scatter([0.18], [y], s=700 if res in {404, 405} else 420, color=color, edgecolors="white", linewidths=1.5, zorder=5)
        ax.text(0.18, y, f"B{res}", ha="center", va="center", color="white", fontsize=11, fontweight="bold")
    for res, y in partner_y.items():
        ax.scatter([0.82], [y], s=380, color=COLORS["target"], edgecolors="white", linewidths=1.2, zorder=5)
        ax.text(0.82, y, f"A{res}", ha="center", va="center", color="white", fontsize=10, fontweight="bold")
    for target_res, partner_res, occupancy in hotspot_links:
        lw = 1.5 + occupancy / 20.0
        color = COLORS["hotspot"] if target_res in {404, 405} else COLORS["gray"]
        ax.plot([0.24, 0.76], [target_y[target_res], partner_y[partner_res]], color=color, alpha=0.72, linewidth=lw)
        ax.text(0.5, (target_y[target_res] + partner_y[partner_res]) / 2, f"{occupancy:.0f}%", fontsize=9, color=COLORS["ink"], ha="center")
    ax.text(
        0.03,
        0.06,
        "B404 主要与 A40 形成双向主链氢键；B405 主要与 A39 侧链形成稳定相互作用。\n这两类联系的占据率显著高于其他界面位点，因此 404/405 被选作后续短肽设计热点。",
        fontsize=10.6,
        color=COLORS["ink"],
        bbox={"facecolor": "white", "edgecolor": COLORS["navy"], "alpha": 0.92, "boxstyle": "round,pad=0.35"},
    )
    ax.set_title("热点级别的界面网络：404/405 是最突出的稳定结合中心", fontsize=17, fontweight="bold", color=COLORS["navy"], pad=14)
    return save_fig(fig, "09_hotspot_network.png")


def create_design_funnel_figure(summary: dict[str, object]) -> Path:
    steps = [
        ("生成候选", summary["generated"]),
        ("覆盖 404/405", summary["hotspot_covered"]),
        ("binder pLDDT ≥ 0.58", summary["binder_confident"]),
        ("iPTM ≥ 0.20", summary["iptm_good"]),
        ("最终通过", summary["final_pass"]),
    ]
    fig, ax = plt.subplots(figsize=(8.5, 5.6))
    widths = np.array([value for _, value in steps], dtype=float)
    widths = widths / widths.max()
    y_positions = np.arange(len(steps))[::-1]
    colors = [COLORS["navy"], COLORS["target"], COLORS["accent"], COLORS["gold"], COLORS["muted_red"]]
    for idx, ((label, value), y) in enumerate(zip(steps, y_positions)):
        width = 0.15 + 0.75 * widths[idx]
        left = 0.5 - width / 2
        rect = plt.Rectangle((left, y - 0.35), width, 0.7, facecolor=colors[idx], alpha=0.9)
        ax.add_patch(rect)
        ax.text(0.5, y, f"{label}\n{value}", ha="center", va="center", fontsize=12, color="white", fontweight="bold")
    ax.set_xlim(0, 1)
    ax.set_ylim(-0.7, len(steps) - 0.3)
    ax.axis("off")
    ax.set_title("设计筛选漏斗：主要瓶颈来自短肽自身置信度，而不是热点覆盖", fontsize=16, fontweight="bold", color=COLORS["navy"], pad=16)
    return save_fig(fig, "10_design_funnel.png")


def create_design_scatter_figure(df: pd.DataFrame) -> Path:
    fig, axes = plt.subplots(1, 2, figsize=(12.5, 5.6))
    scatter = axes[0].scatter(
        df["binder_plddt"],
        df["i_ptm"],
        c=df["best_interface_pae_A"],
        cmap="viridis_r",
        s=42,
        alpha=0.82,
        edgecolors="white",
        linewidths=0.3,
    )
    passed = df[df["secondary_pass"] == "yes"]
    if not passed.empty:
        top = passed.iloc[0]
        axes[0].scatter([top["binder_plddt"]], [top["i_ptm"]], s=150, color=COLORS["muted_red"], edgecolors="black", linewidths=1.2)
        axes[0].annotate(
            "最终候选",
            (top["binder_plddt"], top["i_ptm"]),
            xytext=(top["binder_plddt"] - 0.12, top["i_ptm"] - 0.08),
            arrowprops=dict(arrowstyle="->", color=COLORS["muted_red"], lw=1.5),
            fontsize=11,
            fontweight="bold",
            color=COLORS["muted_red"],
        )
    axes[0].axvline(0.58, color=COLORS["gray"], linestyle="--", linewidth=1.3)
    axes[0].axhline(0.20, color=COLORS["gray"], linestyle="--", linewidth=1.3)
    axes[0].set_xlabel("Binder pLDDT")
    axes[0].set_ylabel("iPTM")
    axes[0].set_title("512 个候选的 AF2 复评分景观", fontsize=13.5)
    style_axis(axes[0])
    cbar = fig.colorbar(scatter, ax=axes[0], shrink=0.88)
    cbar.set_label("Best interface PAE (Å)")

    length_order = sorted(df["binder_length"].unique())
    sns.boxplot(data=df, x="binder_length", y="binder_plddt", order=length_order, ax=axes[1], color=COLORS["gold"])
    sns.stripplot(data=df, x="binder_length", y="binder_plddt", order=length_order, ax=axes[1], color=COLORS["navy"], size=2.5, alpha=0.28)
    axes[1].axhline(0.58, color=COLORS["gray"], linestyle="--", linewidth=1.3)
    axes[1].set_xlabel("Binder length (aa)")
    axes[1].set_ylabel("Binder pLDDT")
    axes[1].set_title("不同长度短肽的置信度分布", fontsize=13.5)
    style_axis(axes[1])

    fig.suptitle("Complexa 设计结果：多数候选能覆盖热点，但只有 1 条短肽达到严格置信度门槛", fontsize=18, fontweight="bold", color=COLORS["navy"])
    return save_fig(fig, "11_design_scatter_and_length.png")


def create_design_structure_figure() -> Path:
    data = load_design_projection_data(DESIGN_FINAL_PDB)
    coords_a = coords_from_records(data["A"])
    coords_b = coords_from_records(data["B"])
    basis, center = build_projection_basis(np.vstack([coords_a, coords_b]))
    proj_a = project_points(coords_a, basis, center)
    proj_b = project_points(coords_b, basis, center)
    fig, ax = plt.subplots(figsize=(8.5, 7.2))
    plot_chain_trace(ax, proj_a, "#C9CDD2", 2.8)
    plot_chain_trace(ax, proj_b, COLORS["accent"], 4.0)
    for resid, label in zip(HOTSPOT_DESIGN, HOTSPOT_DESIGN_LABELS):
        idx = next(i for i, record in enumerate(data["A"]) if record["resid"] == resid)
        xy = proj_a[idx]
        ax.scatter([xy[0]], [xy[1]], s=170, color=COLORS["hotspot"], edgecolors="white", linewidths=1.2, zorder=5)
        ax.text(xy[0] + 0.25, xy[1] + 0.12, label, fontsize=11, color=COLORS["navy"], fontweight="bold")
    for resid in [1, 2, 3, 8, 11]:
        idx = next(i for i, record in enumerate(data["B"]) if record["resid"] == resid)
        record = data["B"][idx]
        xy = proj_b[idx]
        aa = AA3_TO_1.get(record["resname"], "X")
        ax.scatter([xy[0]], [xy[1]], s=90, color=COLORS["accent"], edgecolors="white", linewidths=0.8, zorder=5)
        ax.text(xy[0] + 0.14, xy[1] + 0.1, f"{aa}{resid}", fontsize=10, color=COLORS["accent"], fontweight="bold")
    ax.text(
        0.02,
        0.05,
        "最终候选序列：SMEKIDDLIKR（11 aa）\n热点附近主要接触残基：M2 / E3 / I5 / L8",
        transform=ax.transAxes,
        fontsize=10.6,
        color=COLORS["ink"],
        bbox={"facecolor": "white", "edgecolor": COLORS["accent"], "alpha": 0.92, "boxstyle": "round,pad=0.35"},
    )
    ax.set_title("最终短肽候选的预测复合物构象", fontsize=17, fontweight="bold", color=COLORS["navy"])
    ax.set_xticks([])
    ax.set_yticks([])
    for spine in ax.spines.values():
        spine.set_visible(False)
    return save_fig(fig, "12_final_candidate_structure.png")


def create_metric_card_figure(final_row: pd.Series) -> Path:
    fig, ax = plt.subplots(figsize=(7.4, 4.8))
    ax.axis("off")
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    card = plt.Rectangle((0.03, 0.06), 0.94, 0.88, facecolor="white", edgecolor=COLORS["navy"], linewidth=2.2)
    ax.add_patch(card)
    ax.text(0.07, 0.87, "最终候选关键指标", fontsize=18, fontweight="bold", color=COLORS["navy"])
    ax.text(0.07, 0.78, f"Design ID: {final_row['design_name']}", fontsize=12, color=COLORS["ink"])
    ax.text(0.07, 0.71, "Sequence: SMEKIDDLIKR", fontsize=12.5, color=COLORS["accent"], fontweight="bold")
    metrics = [
        ("Length", f"{int(final_row['binder_length'])} aa"),
        ("Binder pLDDT", f"{final_row['binder_plddt']:.3f}"),
        ("iPTM", f"{final_row['i_ptm']:.3f}"),
        ("Mean interface PAE", f"{final_row['mean_interface_pae_A']:.2f} Å"),
        ("Best interface PAE", f"{final_row['best_interface_pae_A']:.2f} Å"),
        ("Hotspot coverage", f"{int(final_row['hotspot_coverage'])}/2"),
        ("Contact binder residues", f"{int(final_row['contact_binder_residues'])}"),
        ("Min hotspot distance", f"{final_row['min_hotspot_distance']:.2f} Å"),
    ]
    left_x, right_x = 0.08, 0.55
    y = 0.6
    for idx, (label, value) in enumerate(metrics):
        x = left_x if idx < 4 else right_x
        row_y = y - 0.12 * (idx if idx < 4 else idx - 4)
        ax.text(x, row_y, label, fontsize=11, color=COLORS["slate"])
        ax.text(x + 0.22, row_y, value, fontsize=12, color=COLORS["ink"], fontweight="bold")
    ax.text(
        0.08,
        0.12,
        "解释：该候选同时覆盖 404/405，两处热点的最佳界面误差约 2.10 Å，\n是本轮 512 个 8-12 aa 候选中唯一通过严格二筛门槛的短肽。",
        fontsize=10.6,
        color=COLORS["ink"],
        bbox={"facecolor": COLORS["bg"], "edgecolor": COLORS["grid"], "boxstyle": "round,pad=0.35"},
    )
    return save_fig(fig, "13_final_candidate_metrics.png")


def rgb_color(value: str) -> RGBColor:
    named = {
        "white": "FFFFFF",
        "black": "000000",
    }
    token = named.get(value.lower(), value).lstrip("#")
    return RGBColor.from_string(token)


def add_textbox(slide, left, top, width, height, text, font_size=24, bold=False, color=COLORS["ink"], align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = FONT_FAMILY
    run.font.color.rgb = rgb_color(color)
    return box


def add_bullets(slide, left, top, width, height, items, font_size=20, color=COLORS["ink"]):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.name = FONT_FAMILY
        p.font.color.rgb = rgb_color(color)
        p.space_after = Pt(10)
    return box


def apply_slide_style(slide, title: str, subtitle: str | None = None) -> None:
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = rgb_color(COLORS["bg"])
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, PPT_WIDTH, Inches(0.62))
    band.fill.solid()
    band.fill.fore_color.rgb = rgb_color(COLORS["navy"])
    band.line.fill.background()
    add_textbox(slide, Inches(0.45), Inches(0.1), Inches(8.5), Inches(0.35), title, font_size=24, bold=True, color="white")
    if subtitle:
        add_textbox(slide, Inches(9.3), Inches(0.12), Inches(3.6), Inches(0.28), subtitle, font_size=10, color="#DDE6EF", align=PP_ALIGN.RIGHT)


def add_image_fit(slide, path: Path, left: float, top: float, width: float | None = None, height: float | None = None):
    if width is not None and height is not None:
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)
    elif width is not None:
        slide.shapes.add_picture(str(path), left, top, width=width)
    elif height is not None:
        slide.shapes.add_picture(str(path), left, top, height=height)
    else:
        slide.shapes.add_picture(str(path), left, top)


def build_presentation(assets: dict[str, Path], stats: dict[str, dict[str, float]], design_summary: dict[str, object], final_row: pd.Series) -> None:
    prs = Presentation()
    prs.slide_width = PPT_WIDTH
    prs.slide_height = PPT_HEIGHT
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = RGBColor.from_string(COLORS["navy"].lstrip("#"))
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.32))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor.from_string(COLORS["hotspot"].lstrip("#"))
    accent.line.fill.background()
    add_textbox(slide, Inches(0.55), Inches(0.7), Inches(8.8), Inches(0.7), "基于 MD 与 Proteina-Complexa 的靶蛋白热点多肽设计分析", font_size=26, bold=True, color="white")
    add_textbox(slide, Inches(0.58), Inches(1.45), Inches(4.5), Inches(0.5), "汇报版本 | 自动重分析生成", font_size=14, color="#DDE6EF")
    add_textbox(
        slide,
        Inches(0.58),
        Inches(2.1),
        Inches(5.2),
        Inches(1.9),
        "核心主线\n1. MD 证明界面稳定存在\n2. 氢键网络锁定 Chain B 404/405 为关键热点\n3. Complexa 在 8-12 aa 约束下给出 512 个候选\n4. 严格二筛后保留 1 个高可信短肽候选",
        font_size=18,
        color="white",
    )
    add_image_fit(slide, assets["workflow"], Inches(6.6), Inches(0.9), width=Inches(6.1))
    add_textbox(
        slide,
        Inches(6.75),
        Inches(5.9),
        Inches(5.8),
        Inches(0.7),
        "注意：当前 newMD 目录中可复现实验数据为 50 ns 生产轨迹；氢键分析窗口为 15-50 ns。",
        font_size=12,
        color="#F4EBD0",
    )

    slide = prs.slides.add_slide(blank)
    apply_slide_style(slide, "研究设计与证据链", "MD -> hotspot -> peptide design")
    add_image_fit(slide, assets["workflow"], Inches(0.45), Inches(0.9), width=Inches(12.35))
    add_bullets(
        slide,
        Inches(0.7),
        Inches(5.3),
        Inches(12.0),
        Inches(1.3),
        [
            "当前可直接复现的数据长度为 50 ns，因此本次图表与统计全部基于 0-50 ns 重分析。",
            f"氢键分析采用现有结果的 {ANALYSIS_START_NS:.0f}-{ANALYSIS_END_NS:.0f} ns 窗口，以避免把前期未充分稳定段混入热点判断。",
            "汇报逻辑从“体系是否稳定”开始，再推进到“为何选 404/405”，最后落到“8-12 aa 候选是否可靠”。",
        ],
        font_size=16,
    )

    slide = prs.slides.add_slide(blank)
    apply_slide_style(slide, "MD 全程稳定性", "0-50 ns")
    add_image_fit(slide, assets["rmsd"], Inches(0.45), Inches(0.9), width=Inches(8.0))
    add_bullets(
        slide,
        Inches(8.7),
        Inches(1.15),
        Inches(4.1),
        Inches(4.7),
        [
            f"Complex RMSD 在 {ANALYSIS_START_NS:.0f}-{ANALYSIS_END_NS:.0f} ns 的均值为 {stats['rmsd_complex']['mean']:.3f} nm，标准差 {stats['rmsd_complex']['std']:.3f} nm。",
            f"Chain A RMSD 仅 {stats['rmsd_a']['mean']:.3f} ± {stats['rmsd_a']['std']:.3f} nm，说明靶蛋白骨架保持稳定。",
            f"Chain B RMSD 为 {stats['rmsd_b']['mean']:.3f} ± {stats['rmsd_b']['std']:.3f} nm，更像是结合后的构象重排，而非整体解离。",
            "因此后续界面分析不应只看整体 RMSD，而要结合距离、接触数和氢键网络共同判断。",
        ],
        font_size=17,
    )

    slide = prs.slides.add_slide(blank)
    apply_slide_style(slide, "界面是否持续存在", "距离 / 接触数 / 氢键")
    add_image_fit(slide, assets["interface"], Inches(0.45), Inches(0.9), width=Inches(8.0))
    add_bullets(
        slide,
        Inches(8.7),
        Inches(1.15),
        Inches(4.1),
        Inches(4.9),
        [
            f"15-50 ns 最小界面距离均值 {stats['mindist']['mean']:.3f} nm，40-50 ns 仍维持在 {stats['mindist_40_50']['mean']:.3f} nm。",
            f"15-50 ns 接触数均值约 {stats['contacts']['mean']:.0f}，后 10 ns 提升到 {stats['contacts_40_50']['mean']:.0f}。",
            f"氢键数均值 {stats['hbonds']['mean']:.1f}，范围 {stats['hbonds']['min']:.0f}-{stats['hbonds']['max']:.0f}，说明界面存在动态重排但未塌陷。",
            "这一步给出的判断是：体系经历重构后形成了可持续识别的结合界面。",
        ],
        font_size=17,
    )

    slide = prs.slides.add_slide(blank)
    apply_slide_style(slide, "整体紧致性与氢键寿命", "Rg + lifetime")
    add_image_fit(slide, assets["rg_lifetime"], Inches(0.45), Inches(0.95), width=Inches(7.8))
    add_bullets(
        slide,
        Inches(8.55),
        Inches(1.2),
        Inches(4.2),
        Inches(4.8),
        [
            f"40-50 ns Rg 为 {stats['rg_40_50']['mean']:.3f} ± {stats['rg_40_50']['std']:.3f} nm，未见持续松散化趋势。",
            "氢键寿命分布显示界面主要由一组中短寿命、可反复形成的相互作用支撑。",
            "这与蛋白-蛋白界面中“稳定结合 + 局部柔性”的常见模式一致，也解释了 Chain B RMSD 偏高但界面未丢失。",
        ],
        font_size=17,
    )

    slide = prs.slides.add_slide(blank)
    apply_slide_style(slide, "局部柔性与热点暴露", "RMSF")
    add_image_fit(slide, assets["rmsf"], Inches(0.45), Inches(0.95), width=Inches(8.0))
    add_bullets(
        slide,
        Inches(8.7),
        Inches(1.2),
        Inches(4.0),
        Inches(4.7),
        [
            "Chain A 柔性整体较低，支持其作为稳定识别框架。",
            "Chain B 的 404/405 位点处于表面可接近区域，并没有被埋藏到不利于配体进入的位置。",
            "这使 404/405 兼具“能够参与稳定结合”和“后续可被短肽瞄准”的双重特征。",
        ],
        font_size=17,
    )

    slide = prs.slides.add_slide(blank)
    apply_slide_style(slide, "氢键网络锁定 404/405", "Top occupancy")
    add_image_fit(slide, assets["hbond"], Inches(0.45), Inches(0.95), width=Inches(7.9))
    add_bullets(
        slide,
        Inches(8.65),
        Inches(1.12),
        Inches(4.1),
        Inches(4.9),
        [
            "Top 氢键中最显著的界面相互作用集中在 B404 与 B405。",
            "B404 与 A40 形成双向主链氢键，是整个界面里占据率最高的稳定连接。",
            "B405 与 A39 形成高占据率相互作用，构成与 B404 并列的第二核心位点。",
            "因此 404/405 不是‘一次偶然接触’，而是 15-50 ns 内反复出现的界面中心。",
        ],
        font_size=17,
    )

    slide = prs.slides.add_slide(blank)
    apply_slide_style(slide, "结构视角下的热点证据", "reference + snapshots + network")
    add_image_fit(slide, assets["md_structure"], Inches(0.42), Inches(0.95), width=Inches(3.85))
    add_image_fit(slide, assets["snapshot_strip"], Inches(4.45), Inches(0.95), width=Inches(8.4))
    add_image_fit(slide, assets["network"], Inches(0.85), Inches(4.45), width=Inches(11.8))
    add_textbox(
        slide,
        Inches(0.75),
        Inches(6.78),
        Inches(12.0),
        Inches(0.35),
        "从结构快照和氢键网络可以看到：整体构象在变，但 404/405 所在表面始终保持可接近，并持续连接到最关键的对侧残基 A39/A40。",
        font_size=15,
        color=COLORS["ink"],
    )

    slide = prs.slides.add_slide(blank)
    apply_slide_style(slide, "Complexa 设计结果总览", "8-12 aa | n=512")
    add_image_fit(slide, assets["funnel"], Inches(0.6), Inches(1.0), width=Inches(4.2))
    add_image_fit(slide, assets["scatter"], Inches(5.0), Inches(0.95), width=Inches(7.8))
    add_bullets(
        slide,
        Inches(0.72),
        Inches(5.2),
        Inches(12.0),
        Inches(1.3),
        [
            f"共生成 {design_summary['generated']} 个 8-12 aa 候选，其中 {design_summary['hotspot_covered']} 个能够同时覆盖 404/405。",
            f"但只有 {design_summary['binder_confident']} 个在 binder pLDDT 上达到较保守门槛，说明真正的瓶颈是短肽自身稳定性，而不是热点命中率。",
            f"最终严格二筛后仅保留 {design_summary['final_pass']} 个候选，这让最终命中虽然少，但可信度反而更高。",
        ],
        font_size=16,
    )

    slide = prs.slides.add_slide(blank)
    apply_slide_style(slide, "最终候选复合物", "secondary-pass = yes")
    add_image_fit(slide, assets["design_structure"], Inches(0.45), Inches(0.95), width=Inches(7.0))
    add_image_fit(slide, assets["metric_card"], Inches(7.85), Inches(1.15), width=Inches(4.85))
    add_textbox(
        slide,
        Inches(0.65),
        Inches(6.3),
        Inches(6.8),
        Inches(0.45),
        "最终候选序列：SMEKIDDLIKR",
        font_size=18,
        bold=True,
        color=COLORS["accent"],
    )
    add_textbox(
        slide,
        Inches(0.65),
        Inches(6.72),
        Inches(12.0),
        Inches(0.35),
        "结构上可见短肽主要利用 N 端与中段残基贴靠 B404/B405 区域，形成一个覆盖双热点的小界面。",
        font_size=15,
    )

    slide = prs.slides.add_slide(blank)
    apply_slide_style(slide, "结论与建议", "take-home messages")
    add_bullets(
        slide,
        Inches(0.8),
        Inches(1.1),
        Inches(12.0),
        Inches(4.6),
        [
            "MD 重分析支持当前体系在 50 ns 内形成稳定结合界面，不能简单因为整体 RMSD 偏高就判定体系不稳定。",
            "氢键占据率与结构视图一致指向 Chain B 404/405 是最值得优先打击的界面热点。",
            "在 8-12 aa 的严格长度约束下，Proteina-Complexa 可以命中热点，但真正高可信的候选非常稀少，反映该任务本身较难。",
            "当前最值得带去下一轮验证的序列是 SMEKIDDLIKR，建议优先做独立对接复核、短程 MD 复检和体外结合实验。",
        ],
        font_size=20,
    )
    summary_box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(1.35), Inches(4.5), Inches(3.8))
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
    summary_box.line.color.rgb = RGBColor.from_string(COLORS["hotspot"].lstrip("#"))
    add_textbox(slide, Inches(8.3), Inches(1.65), Inches(3.8), Inches(0.4), "推荐汇报结论句", font_size=20, bold=True, color=COLORS["navy"])
    add_textbox(
        slide,
        Inches(8.3),
        Inches(2.2),
        Inches(3.7),
        Inches(2.2),
        "基于 MD 的界面稳定性和氢键占据率分析，我们将病毒蛋白表面的 404/405 界定为关键结合热点；进一步在 8-12 aa 约束下进行 hotspot-guided 设计，最终获得 1 条同时覆盖双热点且 AF2 复评分最优的候选短肽 SMEKIDDLIKR。",
        font_size=16,
        color=COLORS["ink"],
    )

    prs.save(str(PPTX_PATH))


def main() -> None:
    ensure_dirs()

    rmsd_complex = read_xvg(MD_DIR / "rmsd_complex.xvg")
    rmsd_a = read_xvg(MD_DIR / "rmsd_chain_A.xvg")
    rmsd_b = read_xvg(MD_DIR / "rmsd_chain_B.xvg")
    rmsf_a = read_xvg(MD_DIR / "rmsf_chainA.xvg", scale_time_ps_to_ns=False)
    rmsf_b = read_xvg(MD_DIR / "rmsf_chainB.xvg", scale_time_ps_to_ns=False)
    rg = read_xvg(MD_DIR / "gyrate.xvg")
    mindist = read_xvg(MD_DIR / "mindist.xvg")
    contacts = read_xvg(MD_DIR / "num_contact.xvg")
    hb_time, hbonds, hb_pairs = read_xvg_two_series(MD_DIR / "hbnum.xvg")
    hb_life_time_ps, hb_life_prob, hb_life_weighted = read_xvg_two_series(MD_DIR / "hbond_life.xvg", scale_time_ps_to_ns=False)

    hbond_records = parse_hbond_records(PLOT_DIR / "hbonds_corrected.pml")
    atom_df, residue_df = summarize_hbond_records(hbond_records)
    design_df = load_design_table()
    design_summary = compute_design_summary(design_df)
    final_candidates = design_df[design_df["secondary_pass"] == "yes"].copy()
    if final_candidates.empty:
        raise RuntimeError("No final candidate found in design table.")
    final_row = final_candidates.sort_values(["binder_plddt", "i_ptm"], ascending=False).iloc[0]

    stats = {
        "rmsd_complex": compute_stats(rmsd_complex, ANALYSIS_START_NS, ANALYSIS_END_NS),
        "rmsd_a": compute_stats(rmsd_a, ANALYSIS_START_NS, ANALYSIS_END_NS),
        "rmsd_b": compute_stats(rmsd_b, ANALYSIS_START_NS, ANALYSIS_END_NS),
        "mindist": compute_stats(mindist, ANALYSIS_START_NS, ANALYSIS_END_NS),
        "mindist_40_50": compute_stats(mindist, 40.0, 50.0),
        "contacts": compute_stats(contacts, ANALYSIS_START_NS, ANALYSIS_END_NS),
        "contacts_40_50": compute_stats(contacts, 40.0, 50.0),
        "rg_40_50": compute_stats(rg, 40.0, 50.0),
        "hbonds": {
            "mean": float(np.mean(hbonds)),
            "std": float(np.std(hbonds)),
            "min": float(np.min(hbonds)),
            "max": float(np.max(hbonds)),
        },
    }

    assets = {
        "workflow": create_workflow_figure(),
        "rmsd": create_rmsd_figure(rmsd_complex, rmsd_a, rmsd_b),
        "interface": create_interface_persistence_figure(mindist, contacts, hb_time, hbonds),
        "rg_lifetime": create_rg_lifetime_figure(rg, hb_life_time_ps, hb_life_prob, hb_life_weighted),
        "rmsf": create_rmsf_figure(rmsf_a, rmsf_b),
        "hbond": create_hbond_occupancy_figure(atom_df, residue_df),
        "md_structure": create_md_structure_figure(),
        "snapshot_strip": create_snapshot_strip_figure(),
        "network": create_hotspot_network_figure(hbond_records),
        "funnel": create_design_funnel_figure(design_summary),
        "scatter": create_design_scatter_figure(design_df),
        "design_structure": create_design_structure_figure(),
        "metric_card": create_metric_card_figure(final_row),
    }

    build_presentation(assets, stats, design_summary, final_row)

    summary_lines = [
        "LY MD + Complexa report summary",
        f"Data root: {ROOT}",
        f"MD trajectory analyzed: 0-50 ns (reproducible trajectory from {MD_DIR})",
        f"H-bond analysis window: {ANALYSIS_START_NS:.0f}-{ANALYSIS_END_NS:.0f} ns",
        f"Complex RMSD ({ANALYSIS_START_NS:.0f}-{ANALYSIS_END_NS:.0f} ns): {stats['rmsd_complex']['mean']:.4f} ± {stats['rmsd_complex']['std']:.4f} nm",
        f"Chain A RMSD ({ANALYSIS_START_NS:.0f}-{ANALYSIS_END_NS:.0f} ns): {stats['rmsd_a']['mean']:.4f} ± {stats['rmsd_a']['std']:.4f} nm",
        f"Chain B RMSD ({ANALYSIS_START_NS:.0f}-{ANALYSIS_END_NS:.0f} ns): {stats['rmsd_b']['mean']:.4f} ± {stats['rmsd_b']['std']:.4f} nm",
        f"Min interface distance ({ANALYSIS_START_NS:.0f}-{ANALYSIS_END_NS:.0f} ns): {stats['mindist']['mean']:.4f} ± {stats['mindist']['std']:.4f} nm",
        f"Contacts ({ANALYSIS_START_NS:.0f}-{ANALYSIS_END_NS:.0f} ns): {stats['contacts']['mean']:.2f} ± {stats['contacts']['std']:.2f}",
        f"H-bonds ({ANALYSIS_START_NS:.0f}-{ANALYSIS_END_NS:.0f} ns): {stats['hbonds']['mean']:.2f} ± {stats['hbonds']['std']:.2f}",
        f"Top hotspot residues from H-bond occupancy: B{int(residue_df.iloc[0]['residue'])}, B{int(residue_df.iloc[1]['residue'])}",
        f"Design candidates generated: {design_summary['generated']}",
        f"Dual-hotspot covered: {design_summary['hotspot_covered']}",
        f"Final passing candidates: {design_summary['final_pass']}",
        f"Selected sequence: SMEKIDDLIKR",
        f"Selected binder pLDDT: {float(final_row['binder_plddt']):.4f}",
        f"Selected iPTM: {float(final_row['i_ptm']):.4f}",
        f"Selected mean interface PAE: {float(final_row['mean_interface_pae_A']):.4f} A",
        f"Selected best interface PAE: {float(final_row['best_interface_pae_A']):.4f} A",
        f"PPTX output: {PPTX_PATH}",
        f"Asset directory: {ASSET_DIR}",
    ]
    SUMMARY_PATH.write_text("\n".join(summary_lines), encoding="utf-8")


if __name__ == "__main__":
    main()
